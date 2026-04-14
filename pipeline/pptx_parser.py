"""
pipeline/pptx_parser.py
-----------------------
PPTX → List[{page_num, title, text, char_count}]

Mỗi slide = 1 page. Title = text box đầu tiên hoặc shape có type TITLE.

Install: pip install python-pptx
"""

from pathlib import Path


MIN_CHARS = 30   # slide có thể ngắn hơn PDF page


def parse_pptx(file_path: str) -> list[dict]:
    """
    Parse PPTX → List[page_dict] tương thích với pdf_parser output.
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.enum.shapes import PP_PLACEHOLDER
    except ImportError:
        raise ImportError("Chua cai python-pptx. Chay: pip install python-pptx")

    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File khong ton tai: {file_path}")

    prs    = Presentation(file_path)
    pages  = []

    for slide_idx, slide in enumerate(prs.slides):
        page_num = slide_idx + 1
        title    = None
        lines    = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            shape_text = shape.text_frame.text.strip()
            if not shape_text:
                continue

            # Detect title: placeholder idx=0 là title
            is_title = False
            try:
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    is_title = True
            except (ValueError, AttributeError):
                pass

            if is_title and title is None:
                title = shape_text
            else:
                # Body text: lấy từng paragraph
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        lines.append(para_text)

        # Fallback title nếu không detect được
        if title is None and lines:
            title = lines[0]
            lines = lines[1:]
        if title is None:
            title = f"Slide {page_num}"

        text       = "\n".join(lines).strip()
        full_text  = f"{title}\n{text}".strip() if text else title
        char_count = len(full_text)

        if char_count < MIN_CHARS:
            _log(f"  slide {page_num}: SKIP (only {char_count} chars)")
            continue

        pages.append({
            "page_num":   page_num,
            "title":      title,
            "text":       full_text,
            "char_count": char_count,
        })

    return pages


def get_pptx_metadata(file_path: str) -> dict:
    """Metadata cơ bản của PPTX."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        return {
            "total_pages": len(prs.slides),
            "title":       "",
            "author":      "",
            "file_size_kb": round(Path(file_path).stat().st_size / 1024, 1),
        }
    except Exception:
        return {"total_pages": 0, "title": "", "author": "", "file_size_kb": 0}


def _log(msg: str):
    try:
        print(msg)
    except UnicodeEncodeError:
        print(msg.encode("ascii", errors="replace").decode("ascii"))


# ── CLI test ───────────────────────────────────────────────────────
if __name__ == "__main__":
    import sys
    if sys.platform == "win32":
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")

    if len(sys.argv) < 2:
        print("Usage: python -m pipeline.pptx_parser <file.pptx>")
        sys.exit(1)

    pages = parse_pptx(sys.argv[1])
    meta  = get_pptx_metadata(sys.argv[1])

    print(f"\n{'='*55}")
    print(f"  File  : {sys.argv[1]}")
    print(f"  Slides: {meta['total_pages']} total → {len(pages)} kept")
    print(f"{'='*55}\n")

    print(f"{'NUM':<5} {'TITLE':<45} {'CHARS'}")
    print("-" * 60)
    for p in pages[:10]:
        print(f"  {p['page_num']:<3} {p['title'][:43]:<45} {p['char_count']}")
    if len(pages) > 10:
        print(f"  ... ({len(pages)-10} more)")

    if pages:
        print(f"\n--- Sample text (slide 1) ---")
        print(pages[0]["text"][:300])

    # Auto checks
    print("\n--- Auto checks ---")
    assert len(pages) > 0, "FAIL: Khong parse duoc slide nao"
    avg = sum(p["char_count"] for p in pages) / len(pages)
    print(f"  Avg chars/slide: {avg:.0f} {'[OK]' if avg > 50 else '[WARN]'}")
    unique = len(set(p["title"] for p in pages)) / len(pages)
    print(f"  Unique title ratio: {unique:.0%} {'[OK]' if unique > 0.4 else '[WARN] nhieu title trung'}")